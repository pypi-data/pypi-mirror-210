import logging
from pathlib import Path
from typing import Iterable, Optional, cast

import daiquiri
import typer
from pptx.slide import Slide
from typer import FileBinaryRead, FileBinaryWrite, Typer

from dbnomics_pptx_tools import tables
from dbnomics_pptx_tools.cache import SeriesCache
from dbnomics_pptx_tools.cli.cli_utils import (
    load_presentation_metadata,
    open_presentation,
    parse_slide_option,
    parse_slides_option,
)
from dbnomics_pptx_tools.metadata import PresentationMetadata
from dbnomics_pptx_tools.repo import SeriesLoadError, SeriesRepo
from dbnomics_pptx_tools.slides import (
    delete_other_slides,
    extract_slide_id_from_slide_notes,
    find_slide_by_id,
    find_table_shape,
    update_slides,
)
from dbnomics_pptx_tools.tables import format_table

app = Typer()

logger = daiquiri.getLogger(__name__)


DBNOMICS_API_CACHE_DIR_NAME = "dbnomics_api_cache"


@app.callback(context_settings={"help_option_names": ["-h", "--help"]})
def main(verbose: bool = typer.Option(False, "-v", help="Show debug log messages")) -> None:
    """
    DBnomics PowerPoint (pptx) tools.
    """
    daiquiri.setup()
    if verbose:
        daiquiri.set_default_log_levels([("dbnomics_pptx_tools", logging.DEBUG)])


@app.command()
def extract_table_zones(
    input_pptx_file: FileBinaryRead,
    slide_expr: str = typer.Argument(..., help="Slide ID or number"),
    table_name: str = typer.Argument(...),
    metadata_file: Path = typer.Option(..., exists=True, readable=True),
) -> None:
    prs = open_presentation(input_pptx_file)
    prs_slide_ids = [extract_slide_id_from_slide_notes(slide) for slide in cast(Iterable[Slide], prs.slides)]
    slide_number = parse_slide_option(slide_expr, slide_ids=prs_slide_ids)
    slide_id = prs_slide_ids[slide_number - 1]
    if slide_id is None:
        raise typer.BadParameter(f"Slide number {slide_number} does not have an ID defined in the slide notes")

    presentation_metadata = load_presentation_metadata(metadata_file)
    slide_metadata = presentation_metadata.slides.get(slide_id)
    table_spec = None if slide_metadata is None else slide_metadata.tables.get(table_name)

    slide = find_slide_by_id(prs, slide_id=slide_id)
    if slide is None:
        typer.echo(f"Could not find slide wiht ID {slide_id!r}")
        raise typer.Exit(1)
    logger.debug("Found slide %r", slide)

    table_shape = find_table_shape(slide, table_name)
    if table_shape is None:
        typer.echo(f"Could not find table {table_name!r} in slide with ID {slide_id!r}")
        raise typer.Exit(1)
    logger.debug("Found table shape %r", table_shape)

    table = table_shape.table
    logger.debug("Showing table preview:\n%s", format_table(table))

    table_zones = tables.extract_table_zones(table, table_spec=table_spec)
    if table_zones is None:
        typer.echo(f"Could not extract the zones of table {table_name!r}")
        raise typer.Exit(1)

    typer.echo(table_zones)


@app.command()
def fetch(
    metadata_file: Path = typer.Argument(..., exists=True, readable=True),
    dbnomics_api_cache_dir: Path = typer.Option(DBNOMICS_API_CACHE_DIR_NAME),
    skip_existing: bool = typer.Option(False, help="Do not fetch the series that are already stored in the cache."),
) -> None:
    presentation_metadata = load_presentation_metadata(metadata_file)
    series_ids = sorted(presentation_metadata.find_fetchable_series_ids())

    cache = SeriesCache(cache_dir=dbnomics_api_cache_dir)
    repo = SeriesRepo(auto_fetch=True, cache=cache, force=not skip_existing)

    logger.debug("Fetching all the series needed for the presentation: %r...", series_ids)

    fetched_series_ids = []
    skipped_series_ids = []
    for series_id in series_ids:
        if cache.has(series_id) and skip_existing:
            logger.debug("Series %r is already stored in the cache, skipping it", series_id)
            skipped_series_ids.append(series_id)
            continue
        repo.load(series_id)
        fetched_series_ids.append(series_id)

    logger.info(
        "Fetched %d series from DBnomics API, skipped %d series that were already stored in the cache",
        len(fetched_series_ids),
        len(skipped_series_ids),
    )


@app.command()
def show_presentation_metadata_json_schema() -> None:
    """
    Show JSON schema of presentation metadata.
    """

    typer.echo(PresentationMetadata.schema_json(indent=2))


@app.command()
def update(
    input_pptx_file: FileBinaryRead,
    output_pptx_file: FileBinaryWrite,
    auto_fetch: bool = typer.Option(
        True, help="Fetch series when it is not found in the cache, then add it to the cache."
    ),
    dbnomics_api_cache_dir: Path = typer.Option(DBNOMICS_API_CACHE_DIR_NAME),
    fail_fast: bool = typer.Option(False, help="Raise exception when failing to update a slide."),
    force: bool = typer.Option(False, help="Fetch a series even if it is already stored in the cache."),
    metadata_file: Path = typer.Option(..., exists=True, readable=True),
    only_slides_expr: Optional[str] = typer.Option(None, "--slides"),
    save_processed_slides_only: bool = False,
) -> None:
    """
    Update DBnomics data in a PowerPoint (pptx) presentation.
    """
    prs = open_presentation(input_pptx_file)

    only_slides = None
    if only_slides_expr is not None:
        logger.debug("Will process slides %s", only_slides_expr)
        prs_slide_ids = [extract_slide_id_from_slide_notes(slide) for slide in cast(Iterable[Slide], prs.slides)]
        only_slides = parse_slides_option(only_slides_expr, slide_ids=prs_slide_ids)

    if save_processed_slides_only and only_slides is None:
        raise typer.BadParameter("--save-processed-slides-only must be used with --slides")

    presentation_metadata = load_presentation_metadata(metadata_file)

    cache = SeriesCache(cache_dir=dbnomics_api_cache_dir)
    repo = SeriesRepo(auto_fetch=auto_fetch, cache=cache, force=force)

    try:
        update_slides(
            prs, fail_fast=fail_fast, only_slides=only_slides, presentation_metadata=presentation_metadata, repo=repo
        )
    except SeriesLoadError as exc:
        typer.echo(f'{str(exc)} Hint: use the --auto-fetch option or run the "fetch" command first.')
        raise typer.Exit(1)

    if save_processed_slides_only:
        assert only_slides is not None
        delete_other_slides(prs, only_slides=only_slides)

    prs.save(output_pptx_file)
    logger.info("Output presentation was saved as %r", str(output_pptx_file.name))


if __name__ == "__main__":
    app()
