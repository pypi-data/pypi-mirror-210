# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""
import pymysql
import pymssql
import cx_Oracle

import pandas as pd
import numpy as np
from datetime import datetime
import time

from pptx import Presentation
from pptx.util import Pt,Cm
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement

"""
try:
    cx_Oracle.init_oracle_client(lib_dir="D:/programs/datafetch/instantclient_21_7")
except:
    pass
"""
###连接SQL数据库
sun_conn = pymssql.connect('192.168.24.71:2433', 'amperead', 'ampe2017read', 'suntime')
sun_cursor = sun_conn.cursor(as_dict=True)

dsn = cx_Oracle.makedsn('192.168.142.83', '1521', 'VAN')
wind_conn = cx_Oracle.connect(user='RSWINDDB', password='Abc123', dsn=dsn)
wind_cursor = wind_conn.cursor()

#导入并读取info,week,ppt模板表
info = pd.read_excel("info.xlsx",sheet_name='basic')
info = info.astype({'公司名称':'string'})
info_content = pd.read_excel("info.xlsx",sheet_name='content')
df_week = pd.read_excel('tradeweek.xlsx')
prs = Presentation('template43.pptx')

#输入需要制作PPT的公司名
fund_num = input("请输入公司编号（若有多个，请用,分隔）：").split(',')

with open("info.txt", "r") as p:
    s = ",".join(p.readlines()).replace("\n", "")
    p.close()

#fund_num = s.split(",")

#end_date = input("请输入截止日期,格式为yyyy-mm-dd：")
#last_year_date = input("请输入去年年末日期,格式为yyyy-mm-dd：")

end_date = "2023-04-30"
last_year_date = "2022-12-30"
df_labor = pd.read_excel("data_local_tem.xlsx", skiprows=1) #wind, meta, labor临时写在这
df_labor_name = pd.read_excel("data_local_tem.xlsx").iloc[:1].set_index("fundname").T.reset_index().set_index("date")
  

#分别从suntime和wind数据库中获取数据
def get_data(benchmark, fund_id, start_date, db_name):
    ###时间区间确定  ###fund
    if db_name == "zyyxdb":  
        ##如果没有起始日期
        if not pd.isna(start_date):
            sql_fund = "select fund_id,fund_name,statistic_date,swanav from t_fund_nv_data_zyyx where fund_id={} and statistic_date >= '{}' ".format(fund_id,start_date+' 00:00:00.000')
            df_fund = pd.read_sql(sql_fund,sun_conn).rename(columns={'statistic_date':'Date','swanav':'nv'})
        
        else:
            sql_fund = "select fund_id,fund_name,statistic_date,swanav from t_fund_nv_data_zyyx where fund_id={}".format(fund_id)
            df_fund = pd.read_sql(sql_fund,sun_conn).rename(columns={'statistic_date':'Date','swanav':'nv'})
        
        sql_bm = "select \
                    A.S_INFO_WINDCODE, B.S_INFO_NAME, A.TRADE_DT, A.S_DQ_CLOSE\
                  from (\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.AIndexEODPrices) \
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ChinaMutualFundBenchmarkEOD)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ASWSIndexEOD)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.GlobalIndexEOD)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CBIndexEODPrices)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CFutureIndexEODPrices)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ThirdPartyIndexEOD)\
                     union all\
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CMFIndexEOD) \
                     union all \
                     (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CHFIndexEOD)\
                 ) A left join (\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.AIndexDescription)\
                     union all\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CBondDescription)\
                     union all\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CFutureIndexDescription)\
                     union all\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CMFIndexDescription)\
                     union all\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CBIndexDescription)\
                     union all\
                     (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CHFIndexDescription)\
                 ) B on A.S_INFO_WINDCODE = B.S_INFO_WINDCODE\
                 where A.S_INFO_WINDCODE ='{}'  \
                 and A.TRADE_DT >= '20100101' ".format(benchmark)
        df_bm = pd.read_sql(sql_bm,wind_conn).rename(columns={'S_INFO_WINDCODE':'benchmark_id','S_INFO_NAME':'benchmark_name','TRADE_DT':'Date','S_DQ_CLOSE':'close_price'})    
        
    else:  ###如何有起始日期
    # elif db_name in ["labor"]:
        df_fund = df_labor.loc[:, ["date", fund_id]].dropna().reset_index(drop=True).\
                      rename({fund_id: "nv", "date": "Date"}, axis=1).\
                      assign(fund_id=fund_id).\
                      assign(fund_name=df_labor_name.loc[fund_id, "index"]).\
                      loc[:, ["fund_id", "fund_name", "Date", "nv"]]        
        df_fund = df_fund.replace(0,np.nan).dropna()
        
    sql_bm = "select A.S_INFO_WINDCODE, B.S_INFO_NAME, A.TRADE_DT, A.S_DQ_CLOSE\
              from (\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.AIndexEODPrices) \
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ChinaMutualFundBenchmarkEOD)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ASWSIndexEOD)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.GlobalIndexEOD)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CBIndexEODPrices)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CFutureIndexEODPrices)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.ThirdPartyIndexEOD)\
                  union all\
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CMFIndexEOD) \
                  union all \
                  (select S_INFO_WINDCODE, TRADE_DT, S_DQ_CLOSE from WINDDF.CHFIndexEOD)\
              ) A left join (\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.AIndexDescription)\
                  union all\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CBondDescription)\
                  union all\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CFutureIndexDescription)\
                  union all\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CMFIndexDescription)\
                  union all\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CBIndexDescription)\
                  union all\
                  (select S_INFO_WINDCODE, S_INFO_NAME from WINDDF.CHFIndexDescription)\
              ) B on A.S_INFO_WINDCODE = B.S_INFO_WINDCODE\
              where A.S_INFO_WINDCODE ='{}'  \
              and A.TRADE_DT >= '20100101'".format(benchmark)
              
    df_bm = pd.read_sql(sql_bm,wind_conn).rename(columns={'S_INFO_WINDCODE':'benchmark_id','S_INFO_NAME':'benchmark_name','TRADE_DT':'Date','S_DQ_CLOSE':'close_price'})     
    return df_fund,df_bm 
 
###数据预处理
###日期格式统一
def date_format(df,field):
    for i in range(df.shape[0]):
        df.iloc[i,1] = str(df[field][i])
    df[field] = pd.to_datetime(df[field])
    
### merge并得到df
def merge_data(df_fund,df_tradeweek,df_benchmark,date,fund_nv,benchmark_close_price):
    start_dt = df_fund[date].iloc[0]
    if end_date == "":
        end_dt = df_fund[date][df_fund.shape[0]-1]
    else:
        end_dt = end_date         
    tradeweek = df_week[(df_week[date]<=end_dt) & (df_week[date]>=start_dt)]
    df_benchmark = df_bm.drop(columns=['benchmark_id','benchmark_name'])    
    df = df_fund.merge(tradeweek, on=date,how = 'right').merge(df_benchmark,on=date,how = 'left')
    ##排序+缺失值+归一化
    df = df.sort_values(by=date,ascending=True).fillna(method='bfill')
    df.iloc[:,3:] = df.iloc[:,3:]/df.iloc[0,3:]   
    ##收益率序列    
    df['product_rtn'] = df[fund_nv]/df[fund_nv].shift(1)-1
    df['bm_rtn'] = df[benchmark_close_price]/df[benchmark_close_price].shift(1)-1
    ##超额（画图+计算）
    df['excess'] = df[fund_nv] - df[benchmark_close_price]+1
    df['ex_rtn'] =  df['product_rtn'] - df['bm_rtn']
    return df

#ytd序列
def ytd_list(df,field,date):
    ytd_df = df[df[field] >= date]
    ytd_df.reset_index(inplace=True,drop=True)
    return ytd_df 

###指标计算
def ann_return(df,field,plus=52):
    ann_return = (df[field][df.shape[0]-1] / df[field][0] )** (plus / df.shape[0])-1
    return ann_return
def ann_vol(df,field,plus=52):
    ann_vol = df[field].std() * np.sqrt(plus)
    return ann_vol       
def maxdrawdown(df,field):
    bm_maxdrawdown = ((df[field]-df[field].cummax()) / df[field].cummax()).min() 
    return bm_maxdrawdown
def ytd_ret(df,field):
    ytd_rtn = df[field][df.shape[0]-1] / ytd_df[field][0] - 1
    return ytd_rtn    
def itd_ret(df,field):
    itd_rtn = df[field][df.shape[0]-1] / df[field][0] - 1    
    return itd_rtn    


def compute_data(df,fund_nv,benchmark_close_price,product_rtn,bm_rtn,excess,ex_rtn):   
    product_ann_rtn = ann_return(df,fund_nv)
    bm_ann_rtn = ann_return(df,benchmark_close_price)            
    ex_ann_rtn = ann_return(df,excess)
    product_itd_rtn = itd_ret(df,fund_nv)
    bm_itd_rtn = itd_ret(df,benchmark_close_price) 
    ex_itd_rtn = product_itd_rtn - bm_itd_rtn
    product_ann_vol = ann_vol(df,product_rtn)
    bm_ann_vol = ann_vol(df,bm_rtn)
    ex_ann_vol = ann_vol(df,ex_rtn)    
    product_ytd_rtn = ytd_ret(df,fund_nv)
    bm_ytd_rtn = ytd_ret(df,benchmark_close_price)
    ex_ytd_rtn = product_ytd_rtn - bm_ytd_rtn
    product_maxdrawdown = maxdrawdown(df,fund_nv)
    bm_maxdrawdown = maxdrawdown(df,benchmark_close_price)
    ex_maxdrawdown = maxdrawdown(df,excess)
    product_sharpe = product_ann_rtn / product_ann_vol
    bm_sharpe = bm_ann_rtn / bm_ann_vol    
    ex_sharpe = ex_ann_rtn / ex_ann_vol
    return product_ann_rtn,bm_ann_rtn,ex_ann_rtn,product_itd_rtn,bm_itd_rtn,ex_itd_rtn,\
        product_ann_vol,bm_ann_vol,ex_ann_vol,product_ytd_rtn,bm_ytd_rtn,ex_ytd_rtn,\
        product_maxdrawdown,bm_maxdrawdown,ex_maxdrawdown,product_sharpe,bm_sharpe,ex_sharpe,
    
#小数变成百分数
def float_to_percent(ratio_data):
    ratio_data = '%.2f%%' % (ratio_data*100 )
    return ratio_data

###PPT部分
def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
        
##管理人介绍表格样式
def company_introdution_title_table(cell,text,font_rgb,fore_rgb):
    cell.text_frame.paragraphs[0].text = text
    cell.text_frame.paragraphs[0].font.name = '华文楷体'
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = font_rgb
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = fore_rgb    
def company_introdution_content_table(cell,text,font_rgb,fore_rgb):
    cell.text_frame.paragraphs[0].text = text
    cell.text_frame.paragraphs[0].font.name = '华文楷体'
    cell.text_frame.paragraphs[0].font.bold = False
    cell.text_frame.paragraphs[0].font.color.rgb = font_rgb
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = fore_rgb
##计算指标表格样式
def compute_table(cell,text,font_rgb,fore_rgb):
    cell.text_frame.paragraphs[0].text = text
    cell.text_frame.paragraphs[0].font.name = '华文楷体'
    cell.text_frame.paragraphs[0].font.bold = False
    cell.text_frame.paragraphs[0].font.color.rgb = font_rgb
    cell.text_frame.paragraphs[0].font.size = Pt(9)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = fore_rgb
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
##标题内容
def text_content(name,content):
    name.text = content
##line_chart设置
def chart_line_1(placeholder_name,category,series_1,data_1,series_2,data_2):
    chart_data = CategoryChartData()
    chart_data.categories = category
    chart_data.add_series(series_1,data_1)
    chart_data.add_series(series_2,data_2)
    chart = placeholder_name.insert_chart(XL_CHART_TYPE.LINE,chart_data).chart
    chart.has_legend = True  
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.name = '华文楷体'
    chart.legend.font.size = Pt(10)
    chart.font.size = Pt(10)
    chart.font.name = '华文楷体'
    
def chart_line_2(placeholder_name,category,series_1,data_1,series_2,data_2,series_3,data_3):
    chart_data = CategoryChartData()
    chart_data.categories = category
    chart_data.add_series(series_1,data_1)
    chart_data.add_series(series_2,data_2)
    chart_data.add_series(series_3,data_3)
    chart = placeholder_name.insert_chart(XL_CHART_TYPE.LINE,chart_data).chart
    chart.has_legend = True  
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.name = '华文楷体'
    chart.legend.font.size = Pt(10)
    chart.font.size = Pt(10)
    chart.font.name = '华文楷体'    
    
def details_cpt():        
    ###Details   
    product_ann_rtn = float_to_percent(compute_results[0][0])
    bm_ann_rtn = float_to_percent(compute_results[0][1])
    ex_ann_rtn = float_to_percent(compute_results[0][2])
    product_itd_rtn = float_to_percent(compute_results[0][3])
    bm_itd_rtn = float_to_percent(compute_results[0][4])
    ex_itd_rtn = float_to_percent(compute_results[0][5])
    product_ann_vol = float_to_percent(compute_results[0][6])
    bm_ann_vol = float_to_percent(compute_results[0][7])
    ex_ann_vol = float_to_percent(compute_results[0][8])
    product_ytd_rtn = float_to_percent(compute_results[0][9])
    bm_ytd_rtn = float_to_percent(compute_results[0][10])
    ex_ytd_rtn = float_to_percent(compute_results[0][11])
    product_maxdrawdown = float_to_percent(compute_results[0][12])
    bm_maxdrawdown = float_to_percent(compute_results[0][13])
    ex_maxdrawdown = float_to_percent(compute_results[0][14])
    product_sharpe = round(compute_results[0][15],2)
    bm_sharpe = round(compute_results[0][16],2)
    ex_sharpe = round(compute_results[0][17],2)
    
    ##按照模板新添一张幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    ##确定占位符及其类型
    for zwf in slide.placeholders:
        info = zwf.placeholder_format
        print(f'索引{info.idx},名称{zwf.name},类型{info.type}')
        zwf.text = f'索引{info.idx},名称{zwf.name},类型{info.type}'
    
    company_intro = slide.placeholders[23] 
    notice_source = slide.placeholders[21] 
    notice_time = slide.placeholders[19]
    company_intro_table = slide.placeholders[22] 
    ratio_table = slide.placeholders[20]    
    line_chart = slide.placeholders[15]
    line_chart_title = slide.placeholders[14]
    title = slide.placeholders[13]
    
    ##table插入及宽高设置
    intro_row = 3
    intro_col = 2
    ratio_col = 7
    if strategy == '指数增强策略':
        ratio_row = 4
        intro = company_intro_table.insert_table(rows=intro_row,cols=intro_col).table
        table = ratio_table.insert_table(rows=ratio_row,cols=ratio_col).table
        intro.columns[0].width = Cm(1)
        intro.columns[1].width = Cm(9.8)    
        intro.rows[0].height = Cm(2.9)    
        intro.rows[1].height = Cm(4.2)
        intro.rows[2].height = Cm(4.9)
        table.columns[0].width = Cm(1.8)
        for i in range(1,ratio_col):
           table.columns[i].width = Cm(1.6)   
        for i in range(ratio_row):
           table.rows[i].height = Cm(0.8)
        ##具体内容
        for i in range(3):
            for j in range(7):
                set_cell_border(table.cell(i,j), border_color="D9D9D9", border_width='12700')
                set_cell_border(table.cell(3,j), border_color="D9D9D9", border_width='12700')
            set_cell_border(intro.cell(i,0), border_color="D9D9D9", border_width='12700')
            set_cell_border(intro.cell(i,1), border_color="D9D9D9", border_width='12700')

        chart_line_2(line_chart,df.Date,'产品净值',df.nv,'{}'.format(bm_name),df.close_price,'超额收益',df.excess)
        company_introdution_title_table(intro.cell(0,0),manager_intro['栏目'].iloc[0],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_title_table(intro.cell(1,0),manager_intro['栏目'].iloc[1],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_title_table(intro.cell(2,0),manager_intro['栏目'].iloc[2],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(0,1),manager_intro['内容'].iloc[0],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(1,1),manager_intro['内容'].iloc[1],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(2,1),manager_intro['内容'].iloc[2],RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(0,0),'',RGBColor(0,0,0),RGBColor(155,53,25))  
        compute_table(table.cell(0,1),'年化收益',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,2),'年化波动',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,3),'夏普比率',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,4),'最大回撤',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,5),'今年以来',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,6),'成立以来',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(1,0),'产品净值',RGBColor(0,0,0),RGBColor(255,255,255))     
        compute_table(table.cell(2,0),'{}'.format(bm_name),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,0),'超额',RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,1),str(product_ann_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,2),str(product_ann_vol),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,3),str(product_sharpe),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,4),str(product_maxdrawdown),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,5),str(product_ytd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,6),str(product_itd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,1),str(bm_ann_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,2),str(bm_ann_vol),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,3),str(bm_sharpe),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,4),str(bm_maxdrawdown),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,5),str(bm_ytd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,6),str(bm_itd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,1),str(ex_ann_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,2),str(ex_ann_vol),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,3),str(ex_sharpe),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,4),str(ex_maxdrawdown),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,5),str(ex_ytd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(3,6),str(ex_itd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
    else:
        ratio_row = 3 
        intro = company_intro_table.insert_table(rows=intro_row,cols=intro_col).table
        table = ratio_table.insert_table(rows=ratio_row,cols=ratio_col).table
        intro.columns[0].width = Cm(1)
        intro.columns[1].width = Cm(9.8)    
        intro.rows[0].height = Cm(2.9)    
        intro.rows[1].height = Cm(4.2)
        intro.rows[2].height = Cm(4.9)  
        table.columns[0].width = Cm(1.8)
        for i in range(1,ratio_col):
            table.columns[i].width = Cm(1.6)   
        for i in range(ratio_row):
            table.rows[i].height = Cm(1.1)    
        ##具体内容
        for i in range(3):
            for j in range(7):
                set_cell_border(table.cell(i,j), border_color="D9D9D9", border_width='12700')
            set_cell_border(intro.cell(i,0), border_color="D9D9D9", border_width='12700')
            set_cell_border(intro.cell(i,1), border_color="D9D9D9", border_width='12700')
            
        chart_line_1(line_chart,df.Date,'产品净值',df.nv,'{}'.format(bm_name),df.close_price)
        company_introdution_title_table(intro.cell(0,0),manager_intro['栏目'].iloc[0],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_title_table(intro.cell(1,0),manager_intro['栏目'].iloc[1],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_title_table(intro.cell(2,0),manager_intro['栏目'].iloc[2],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(0,1),manager_intro['内容'].iloc[0],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(1,1),manager_intro['内容'].iloc[1],RGBColor(0,0,0),RGBColor(255,255,255))
        company_introdution_content_table(intro.cell(2,1),manager_intro['内容'].iloc[2],RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(0,0),'',RGBColor(0,0,0),RGBColor(155,53,25))  
        compute_table(table.cell(0,1),'年化收益',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,2),'年化波动',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,3),'夏普比率',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,4),'最大回撤',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,5),'今年以来',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(0,6),'成立以来',RGBColor(255,255,255),RGBColor(155,53,25))
        compute_table(table.cell(1,0),'产品净值',RGBColor(0,0,0),RGBColor(255,255,255))     
        compute_table(table.cell(2,0),'{}'.format(bm_name),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,1),str(product_ann_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,2),str(product_ann_vol),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,3),str(product_sharpe),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,4),str(product_maxdrawdown),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,5),str(product_ytd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(1,6),str(product_itd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,1),str(bm_ann_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,2),str(bm_ann_vol),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,3),str(bm_sharpe),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,4),str(bm_maxdrawdown),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,5),str(bm_ytd_rtn),RGBColor(0,0,0),RGBColor(255,255,255))
        compute_table(table.cell(2,6),str(bm_itd_rtn),RGBColor(0,0,0),RGBColor(255,255,255)) 

    
    # 转年月日格式(字符串文本)
    df['年月日'] = df['Date'].apply(lambda x: x.strftime('%Y%m%d'))
    df['年']=df['Date'].dt.year  
    df['月']=df['Date'].dt.month
    df['日']=df['Date'].dt.day
    ###日期写数据的最新日期
    
    notice_time_content='资料来源：管理人，中金公司股票业务部。          数据区间：{}/{}/{} - {}/{}/{}'.format(df['年'][0],df['月'][0],df['日'][0],df['年'][df.shape[0]-1],df['月'][df.shape[0]-1],df['日'][df.shape[0]-1])
    notice_source_content='资料来源：管理人，中金公司股票业务部。'
    title_content='管理人（{}）：{}策略（{}）'.format(manager_name,strategy,bm_name)    
    text_content(notice_time, notice_time_content)
    text_content(notice_source, notice_source_content)
    text_content(title,title_content)
    text_content(line_chart_title,"代表产品业绩" )
    text_content(company_intro,'管理人介绍')
   
    prs.save('一页ppt.pptx')
            
for i in range(len(fund_num)):
    fund_num[i]=int(fund_num[i])
    
for i in range(len(fund_num)):
    manager = info[info['编号']==fund_num[i]]
    manager_intro = info_content[info_content['公司名称']==manager['公司名称'].iloc[0]]
    manager_name = manager['公司名称'].iloc[0]
    benchmark = manager['比较基准'].iloc[0]
    fund_id = manager['代码'].iloc[0]
    start_date = manager['起始日期'].iloc[0]
    strategy = manager['代表策略'].iloc[0]
    data_source = manager['数据源'].iloc[0]
    df_fund,df_bm = get_data(benchmark,fund_id,start_date, data_source)
    bm_name = df_bm['benchmark_name'][0]
    date_format(df_bm,'Date')
    df = merge_data(df_fund,df_week,df_bm,'Date','nv','close_price')
    ytd_df = ytd_list(df,'Date',last_year_date)
    compute_dt =  pd.DataFrame(list(compute_data(df,'nv','close_price','product_rtn','bm_rtn','excess','ex_rtn')))
    compute_results = compute_dt.fillna(0)
    details_cpt()
   

